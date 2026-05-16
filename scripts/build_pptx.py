#!/usr/bin/env python3
"""Generates the full course PPTX with proper typography, syntax highlighting, and smart layout."""

from __future__ import annotations

import os
import re
import textwrap
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.mathtext as mathtext
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import Token

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = str(ROOT / "NEXA_Introduction_aux_IA_Generatives.pptx")
OUTPUT = str(ROOT / "NEXA_Computer_Vision_Complet.pptx")
FORMULA_DIR = ROOT / "outputs" / "pptx_formulas"
FORMULA_DIR.mkdir(parents=True, exist_ok=True)

# ── Charte NEXA ──────────────────────────────────────────────
COLOR_BG = RGBColor(0x1E, 0x1E, 0x2E)
COLOR_TITLE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ACCENT = RGBColor(0x4F, 0x81, 0xBD)
COLOR_ACCENT2 = RGBColor(0xC0, 0x50, 0x4D)
COLOR_ACCENT3 = RGBColor(0x9B, 0xBB, 0x59)
COLOR_BODY = RGBColor(0xE0, 0xE0, 0xE0)
COLOR_MUTED = RGBColor(0x88, 0x88, 0x88)
COLOR_CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
COLOR_CODE_BADGE = RGBColor(0x4F, 0x81, 0xBD)
COLOR_BLUE_LIGHT = RGBColor(0x8A, 0xB4, 0xF8)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

# ── Pygments token → RGB mapping (Dracula-inspired on dark bg) ──
TOKEN_COLORS = {
    Token.Keyword: RGBColor(0xFF, 0x79, 0xC6),
    Token.Keyword.Constant: RGBColor(0xFF, 0x79, 0xC6),
    Token.Keyword.Declaration: RGBColor(0xFF, 0x79, 0xC6),
    Token.Keyword.Namespace: RGBColor(0xFF, 0x79, 0xC6),
    Token.Keyword.Reserved: RGBColor(0xFF, 0x79, 0xC6),
    Token.Keyword.Type: RGBColor(0x8B, 0xE9, 0xFD),
    Token.Name.Builtin: RGBColor(0x8B, 0xE9, 0xFD),
    Token.Name.Function: RGBColor(0x50, 0xFA, 0x7B),
    Token.Name.Class: RGBColor(0x50, 0xFA, 0x7B),
    Token.Name.Decorator: RGBColor(0x50, 0xFA, 0x7B),
    Token.Literal.Number: RGBColor(0xBD, 0x93, 0xF9),
    Token.Literal.Number.Integer: RGBColor(0xBD, 0x93, 0xF9),
    Token.Literal.Number.Float: RGBColor(0xBD, 0x93, 0xF9),
    Token.Literal.String: RGBColor(0xF1, 0xFA, 0x8C),
    Token.Literal.String.Doc: RGBColor(0x62, 0x72, 0xA4),
    Token.Comment: RGBColor(0x62, 0x72, 0xA4),
    Token.Comment.Single: RGBColor(0x62, 0x72, 0xA4),
    Token.Comment.Multiline: RGBColor(0x62, 0x72, 0xA4),
    Token.Operator: RGBColor(0xFF, 0x79, 0xC6),
    Token.Operator.Word: RGBColor(0xFF, 0x79, 0xC6),
    Token.Punctuation: RGBColor(0xF8, 0xF8, 0xF2),
    Token.Text: RGBColor(0xF8, 0xF8, 0xF2),
}
FALLBACK_COLOR = RGBColor(0xF8, 0xF8, 0xF2)

# ── Helpers ──────────────────────────────────────────────────

def token_color(ttype) -> RGBColor:
    for cls in ttype if hasattr(ttype, '__iter__') else [ttype]:
        for cls2 in (cls if hasattr(cls, '__iter__') else [cls]) if isinstance(cls, type) else [cls]:
            if cls2 in TOKEN_COLORS:
                return TOKEN_COLORS[cls2]
    return FALLBACK_COLOR


def load_md(path: str) -> str:
    with open(path, encoding="utf-8") as f:
        return f.read()


def img_exists(path: str) -> bool:
    full = ROOT / path if not os.path.isabs(path) else Path(path)
    return full.is_file() and str(full).lower().endswith((".png", ".jpg", ".jpeg"))


def add_textbox(slide, left, top, width, height, word_wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = word_wrap
    return txBox.text_frame


def set_run(run, text, size=14, bold=False, italic=False, color=COLOR_BODY, font_name=FONT_BODY):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def add_para(tf, text, size=14, bold=False, color=COLOR_BODY, font_name=FONT_BODY, spacing_after=4, alignment=None):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(spacing_after)
    if alignment is not None:
        p.alignment = alignment
    run = p.add_run()
    set_run(run, text, size, bold, color=color, font_name=font_name)
    return p


def add_slide_number(slide, num, total):
    """Ajoute un numéro de slide discret en bas à droite."""
    tf = add_textbox(slide, 16.0, 10.2, 3.5, 0.5)
    add_para(tf, f"{num} / {total}", 9, color=COLOR_MUTED, alignment=PP_ALIGN.RIGHT)


def parse_sections(md: str):
    lines = md.split("\n")
    sections = []
    title = ""
    level = 0
    content = []
    in_code = False
    for line in lines:
        if line.startswith("```"):
            in_code = not in_code
            content.append(line)
            continue
        if not in_code and line.startswith("#"):
            if content or title:
                sections.append((title, level, content))
            title = re.sub(r"^#+\s*", "", line).strip()
            level = len(line) - len(line.lstrip("#"))
            content = []
        else:
            content.append(line)
    if content or title:
        sections.append((title, level, content))
    return sections


# ── Slide builders ───────────────────────────────────────────

def build_cover(prs, total_slides):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Détection et reconnaissance\nd'objets"
        elif ph.placeholder_format.idx == 1:
            ph.text = "Module complet — 3 jours"
    tf = add_textbox(slide, 1.5, 5.5, 17.0, 3.0)
    add_para(tf, "OpenCV  ·  CNN  ·  Faster R-CNN  ·  YOLO  ·  Reconnaissance faciale", 18, color=COLOR_ACCENT)
    add_para(tf, "Mastère 1  —  Intelligence Artificielle", 14, color=COLOR_MUTED)
    return slide


def build_toc(prs, total_slides, current):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Sommaire"
    items = [
        ("Jour 1", "Fondamentaux OpenCV, IoU, HOG, SIFT"),
        ("Jour 2", "CNN, Faster R-CNN, métriques de détection"),
        ("Jour 3", "YOLO, comparaison, optimisation du seuil"),
        ("Projet bonus", "Reconnaissance faciale — YuNet, SFace"),
    ]
    tf = add_textbox(slide, 1.5, 2.5, 17.0, 6.0)
    for day, desc in items:
        p = add_para(tf, "", 16, spacing_after=2)
        r = p.add_run()
        set_run(r, f"  {day}  ", 22, bold=True, color=COLOR_ACCENT)
        r = p.add_run()
        set_run(r, desc, 16, color=COLOR_BODY)
    add_slide_number(slide, current, total_slides)
    return slide


def build_section_slide(prs, title: str, subtitle: str = "", total=0, current=0):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
    if subtitle:
        tf = add_textbox(slide, 1.5, 5.0, 17.0, 2.0)
        add_para(tf, subtitle, 20, color=COLOR_ACCENT)
    add_slide_number(slide, current, total)
    return slide


def build_concept_slide(prs, title: str, bullets: list[str], total=0, current=0):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
    tf = add_textbox(slide, 0.8, 1.8, 18.4, 8.0)
    for i, line in enumerate(bullets[:12]):
        clean = re.sub(r"^-\s*", "", line)
        if line.startswith("- "):
            p = add_para(tf, "", 15, spacing_after=3)
            p.level = 0
            r = p.add_run()
            set_run(r, "▸  ", 13, color=COLOR_ACCENT)
            r = p.add_run()
            set_run(r, clean, 15, color=COLOR_BODY)
        elif "**" in line:
            parts = re.split(r"(\*\*.*?\*\*)", line)
            p = add_para(tf, "", 15, spacing_after=3)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    r = p.add_run()
                    set_run(r, part[2:-2], 15, bold=True, color=COLOR_BODY)
                elif part.strip():
                    r = p.add_run()
                    set_run(r, part, 15, color=COLOR_BODY)
        else:
            add_para(tf, line, 15, spacing_after=3, color=COLOR_BODY)
    add_slide_number(slide, current, total)
    return slide


def build_code_slide(prs, title: str, code: str, filename: str = "", total=0, current=0):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    lines = code.split("\n")

    # Fond sombre texte
    tf_label = add_textbox(slide, 0.8, 0.3, 18.0, 0.8)
    r = add_para(tf_label, title, 20, bold=True, color=COLOR_TITLE)

    # Badge fichier
    if filename:
        badge = add_textbox(slide, 15.5, 0.3, 4.0, 0.6)
        bg = badge.paragraphs[0]
        bg.alignment = PP_ALIGN.RIGHT
        r = bg.add_run()
        set_run(r, f"  {filename}  ", 10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), font_name=FONT_CODE)
        # Add background shape for badge
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1, Inches(15.5), Inches(0.3), Inches(3.0), Inches(0.55)
        )  # rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CODE_BADGE
        shape.line.fill.background()
        tf2 = shape.text_frame
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        r2 = tf2.paragraphs[0].add_run()
        set_run(r2, filename, 10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), font_name=FONT_CODE)

    # Code block background
    from pptx.util import Emu
    bg_shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.3), Inches(19.0), Inches(8.8)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    bg_shape.line.fill.background()

    # Colorized code
    tf = add_textbox(slide, 0.8, 1.5, 18.4, 8.3)
    max_lines = 28
    tokens = list(lex(code, PythonLexer()))
    para = tf.paragraphs[0]
    para.space_after = Pt(0)
    line_count = 0
    current_line = ""
    for ttype, value in tokens:
        if "\n" in value:
            parts = value.split("\n")
            for pi, part in enumerate(parts):
                if current_line:
                    r = para.add_run()
                    set_run(r, current_line, 11, font_name=FONT_CODE, color=token_color(ttype))
                    current_line = ""
                if pi < len(parts) - 1:
                    if line_count >= max_lines:
                        break
                    line_count += 1
                    para = tf.add_paragraph()
                    para.space_after = Pt(0)
                    para.space_before = Pt(0)
                else:
                    current_line = part
        else:
            if line_count >= max_lines:
                break
            current_line += value

    if current_line and line_count < max_lines:
        r = para.add_run()
        set_run(r, current_line, 11, font_name=FONT_CODE, color=token_color(ttype))

    add_slide_number(slide, current, total)
    return slide


def build_formula_slide(prs, title: str, formula_tex: str, explanation: str = "", total=0, current=0):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title

    # Render formula as image
    fig, ax = plt.subplots(figsize=(12, 2.5))
    ax.axis("off")
    ax.set_facecolor("none")
    fig.patch.set_alpha(0)
    try:
        ax.text(0.5, 0.5, f"${formula_tex}$", fontsize=28, ha="center", va="center",
                color="white", transform=ax.transAxes)
        img_path = str(FORMULA_DIR / f"formula_{current}.png")
        fig.savefig(img_path, dpi=150, bbox_inches="tight", transparent=True, pad_inches=0.3)
        plt.close(fig)
        if os.path.isfile(img_path):
            slide.shapes.add_picture(img_path, Inches(2), Inches(2.2), Inches(16), Inches(4))
    except Exception:
        plt.close(fig)

    if explanation:
        tf = add_textbox(slide, 1.5, 7.0, 17.0, 2.5)
        add_para(tf, explanation, 14, color=COLOR_BODY)
    add_slide_number(slide, current, total)
    return slide


def build_image_slide(prs, title: str, img_path: str, caption: str = "", total=0, current=0):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
    full = ROOT / img_path if not os.path.isabs(img_path) else Path(img_path)
    if full.is_file():
        try:
            im = Image.open(full)
            w, h = im.size
            scale = min(14.0 / w, 6.5 / h, 1.0)
            slide.shapes.add_picture(str(full), Inches(3.0), Inches(1.8), Inches(w * scale), Inches(h * scale))
        except Exception:
            pass
    if caption:
        tf = add_textbox(slide, 1.5, 8.5, 17.0, 2.0)
        add_para(tf, caption, 13, color=COLOR_MUTED)
    add_slide_number(slide, current, total)
    return slide


def build_comparison_slide(prs, title: str, headers: list[str], rows: list[list[str]], total=0, current=0):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
    tf = add_textbox(slide, 0.5, 1.8, 19.0, 8.0)
    # Header
    header_text = "  │  ".join(headers)
    p = add_para(tf, "", 14, spacing_after=8)
    r = p.add_run()
    set_run(r, header_text, 14, bold=True, color=COLOR_ACCENT)
    for row in rows:
        row_text = "  │  ".join(row)
        add_para(tf, row_text, 13, spacing_after=4, color=COLOR_BODY)
    add_slide_number(slide, current, total)
    return slide


def build_summary_slide(prs, title: str, points: list[str], total=0, current=0):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
    tf = add_textbox(slide, 1.5, 2.5, 17.0, 6.5)
    for pt in points:
        p = add_para(tf, "", 16, spacing_after=8)
        r = p.add_run()
        set_run(r, "✓  ", 16, bold=True, color=COLOR_ACCENT3)
        r = p.add_run()
        set_run(r, pt, 16, color=COLOR_BODY)
    add_slide_number(slide, current, total)
    return slide


# ── Content routing ──────────────────────────────────────────

def render_markdown(prs, md_content, prefix, slide_counter, total_slides):
    """Parse markdown and create slides. Returns the last slide counter value."""
    sections = parse_sections(md_content)
    pending_img = None
    last_section_title = ""

    for title, level, content in sections:
        if not title:
            continue
        raw_text = "\n".join(content)
        clean_text = raw_text
        # Extract images
        imgs = re.findall(r"!\[.*?\]\((.*?)\)", clean_text)
        img_ref = None
        for ref in imgs:
            if img_exists(ref):
                img_ref = ref
                break
        # Extract code blocks
        code_blocks = re.findall(r"```(?:python|bash|text|json)?\n(.*?)```", clean_text, re.DOTALL)

        # Extract math formulas
        formulas = re.findall(r"\$\$(.*?)\$\$", clean_text, re.DOTALL)

        # Extract tables (simple pipe tables)
        tables = re.findall(r"^\|.*\|$", clean_text, re.MULTILINE)

        # Clean text for bullets
        clean_for_bullets = clean_text
        clean_for_bullets = re.sub(r"!\[.*?\]\(.*?\)", "", clean_for_bullets)
        clean_for_bullets = re.sub(r"```.*?```", "", clean_for_bullets, flags=re.DOTALL)
        clean_for_bullets = re.sub(r"\$\$.*?\$\$", "", clean_for_bullets, flags=re.DOTALL)
        clean_for_bullets = re.sub(r"\$.*?\$", "", clean_for_bullets)
        clean_for_bullets = re.sub(r"\|.*\|", "", clean_for_bullets)
        clean_for_bullets = re.sub(r"[-]{3,}", "", clean_for_bullets)
        bullets = [l.strip() for l in clean_for_bullets.split("\n") if l.strip()]

        section_title = f"{prefix} — {title}" if level > 1 else title

        # Level 1: major section
        if level <= 1:
            slide_counter[0] += 1
            build_section_slide(prs, title, total=total_slides, current=slide_counter[0])
            continue

        # Level 2: section header + content
        if level == 2 and bullets:
            slide_counter[0] += 1
            build_section_slide(prs, title, total=total_slides, current=slide_counter[0])
            chunk = []
            for b in bullets:
                chunk.append(b)
                if len(chunk) >= 8:
                    slide_counter[0] += 1
                    build_concept_slide(prs, section_title, chunk, total=total_slides, current=slide_counter[0])
                    chunk = []
            if chunk:
                slide_counter[0] += 1
                build_concept_slide(prs, section_title, chunk, total=total_slides, current=slide_counter[0])
            continue

        # Level 3+:
        if img_ref:
            slide_counter[0] += 1
            build_image_slide(prs, section_title, img_ref, caption=section_title, total=total_slides, current=slide_counter[0])
            continue

        # Extract filename from context for code slides
        filename = ""
        for line in content:
            m = re.match(r"`([^`]+\.py)`", line)
            if m:
                filename = m.group(1)
                break
        if not filename and "lab.py" in raw_text:
            m = re.search(r"([a-zA-Z0-9_]+\.py)", raw_text)
            if m:
                filename = m.group(1)

        if code_blocks:
            for cb in code_blocks:
                lines_cb = cb.split("\n")
                for ci in range(0, len(lines_cb), 28):
                    chunk = "\n".join(lines_cb[ci:ci + 28])
                    slide_counter[0] += 1
                    build_code_slide(prs, section_title, chunk, filename=filename, total=total_slides, current=slide_counter[0])
            continue

        if formulas:
            for form in formulas:
                slide_counter[0] += 1
                build_formula_slide(prs, section_title, form.strip(), total=total_slides, current=slide_counter[0])
            continue

        if bullets:
            chunk = []
            for b in bullets:
                chunk.append(b)
                if len(chunk) >= 8:
                    slide_counter[0] += 1
                    build_concept_slide(prs, section_title, chunk, total=total_slides, current=slide_counter[0])
                    chunk = []
            if chunk:
                slide_counter[0] += 1
                build_concept_slide(prs, section_title, chunk, total=total_slides, current=slide_counter[0])


def main():
    print("Ouverture du template...")
    prs = Presentation(TEMPLATE)

    # Estimate total slides (approximate for numbering)
    TOTAL_ESTIMATE = 350

    slide_counter = [0]

    # Cover
    slide_counter[0] += 1
    build_cover(prs, TOTAL_ESTIMATE)

    # TOC
    slide_counter[0] += 1
    build_toc(prs, TOTAL_ESTIMATE, slide_counter[0])

    # JOUR 1
    print("JOUR-01.md...")
    build_section_slide(prs, "Jour 1", "Fondamentaux de vision par ordinateur\ndescripteurs classiques", TOTAL_ESTIMATE, slide_counter[0])
    slide_counter[0] += 1
    render_markdown(prs, load_md(str(ROOT / "JOUR-01.md")), "Jour 1", slide_counter, TOTAL_ESTIMATE)

    # JOUR 2
    print("JOUR-02.md...")
    slide_counter[0] += 1
    build_section_slide(prs, "Jour 2", "CNN et Faster R-CNN", TOTAL_ESTIMATE, slide_counter[0])
    render_markdown(prs, load_md(str(ROOT / "JOUR-02.md")), "Jour 2", slide_counter, TOTAL_ESTIMATE)

    # JOUR 3
    print("JOUR-03.md...")
    slide_counter[0] += 1
    build_section_slide(prs, "Jour 3", "YOLO, comparaison et optimisation", TOTAL_ESTIMATE, slide_counter[0])
    render_markdown(prs, load_md(str(ROOT / "JOUR-03.md")), "Jour 3", slide_counter, TOTAL_ESTIMATE)

    # PROJET
    print("PROJET-RECONNAISSANCE-FACIALE.md...")
    slide_counter[0] += 1
    build_section_slide(prs, "Projet bonus", "Reconnaissance faciale — Keanu Reeves", TOTAL_ESTIMATE, slide_counter[0])
    render_markdown(prs, load_md(str(ROOT / "PROJET-RECONNAISSANCE-FACIALE.md")), "Projet", slide_counter, TOTAL_ESTIMATE)

    # Save
    actual = slide_counter[0]
    print(f"Sauvegarde ({actual} slides)...")
    prs.save(OUTPUT)
    size_mb = os.path.getsize(OUTPUT) / 1024 / 1024
    print(f"OK → {OUTPUT}  ({size_mb:.1f} MB, {actual} slides)")


if __name__ == "__main__":
    main()
